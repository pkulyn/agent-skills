#!/usr/bin/env python3
"""html2pptx: 将 my-design HTML 幻灯片导出为 PPTX。

双模导出：
  1. 保真模式（默认）：截图→图片嵌入，像素级还原
  2. 可编辑模式：解析 HTML→原生 PPTX 元素，可编辑

用法：
  from html2pptx import export

  # 保真模式（推荐）
  export("slides.html", "output.pptx", mode="fidelity")

  # 可编辑模式
  export("slides.html", "output.pptx", mode="editable")

  # CLI
  python3 html2pptx.py slides.html output.pptx [--mode fidelity|editable]
"""
import json, base64, os, time, re, sys
import urllib.request
from pathlib import Path

# ── CDP 截图 ──

CDP = "http://127.0.0.1:18800"
_ws = None

def _get_ws():
    global _ws
    if _ws is not None:
        try:
            _ws.ping()
            return _ws
        except:
            _ws = None
    import websocket
    targets = json.loads(urllib.request.urlopen(f"{CDP}/json").read())
    page = next(t for t in targets if t.get('type') == 'page')
    _ws = websocket.create_connection(page['webSocketDebuggerUrl'], suppress_origin=True)
    return _ws

def _cdp(method, params=None, retry=2):
    ws = _get_ws()
    for attempt in range(retry):
        try:
            ws.send(json.dumps({"id": 1, "method": method, "params": params or {}}))
            result = json.loads(ws.recv())
            # Skip events, wait for our response
            while 'id' not in result:
                result = json.loads(ws.recv())
            return result
        except Exception as e:
            if attempt < retry - 1:
                global _ws
                _ws = None
                ws = _get_ws()
            else:
                raise

def _screenshot(width=1920, height=1080):
    result = _cdp("Page.captureScreenshot", {
        "format": "png", "quality": 95,
        "clip": {"x": 0, "y": 0, "width": width, "height": height, "scale": 1}
    })
    if 'result' in result and 'data' in result['result']:
        return base64.b64decode(result['result']['data'])
    return None


# ── HTML 解析（可编辑模式） ──

def _parse_slides(html_content):
    """从 HTML 中提取每页幻灯片的内容结构。"""
    slides = []
    
    # Split by section.s boundaries
    # Match both <section class="s ..."> and <section class="s" ...>
    sections = re.split(r'<section\s+class="s[^"]*"[^>]*>', html_content)
    
    for section in sections[1:]:
        slide = {"texts": [], "shapes": [], "notes": ""}
        
        # Extract data-notes attribute if present
        notes_match = re.search(r'data-notes="([^"]*)"', section[:200])
        if notes_match:
            slide["notes"] = notes_match.group(1)
        
        # Find all div elements with style containing font-size
        # Pattern: <div style="...font-size:Npx...">content</div>
        div_pattern = r'<(?:div|span)[^>]*style="[^"]*font-size:\s*(\d+)px[^"]*"[^>]*>(.*?)</(?:div|span)>'
        
        for match in re.finditer(div_pattern, section, re.DOTALL):
            font_size = int(match.group(1))
            raw_content = match.group(2)
            
            # Clean HTML tags from content
            text = re.sub(r'<[^>]+>', '', raw_content).strip()
            if not text or len(text) < 2:
                continue
            
            # Skip very small text (page indicators, etc.)
            if font_size < 14:
                continue
            
            # Determine level by size
            if font_size >= 52:
                level = "title"
            elif font_size >= 36:
                level = "subtitle"
            elif font_size >= 24:
                level = "heading"
            else:
                level = "body"
            
            # Extract color
            color_match = re.search(r'color:\s*(#[A-Fa-f0-9]{3,8}|rgb\([^)]+\))', match.group(0))
            color = color_match.group(1) if color_match else None
            
            # Extract font-weight
            weight_match = re.search(r'font-weight:\s*(\d+)', match.group(0))
            weight = int(weight_match.group(1)) if weight_match else 400
            
            slide["texts"].append({
                "text": text,
                "size": font_size,
                "level": level,
                "color": color,
                "bold": weight >= 600,
            })
        
        # Extract SVG count (charts)
        svg_count = len(re.findall(r'<svg', section))
        if svg_count > 0:
            slide["has_chart"] = True
            slide["chart_count"] = svg_count
        
        slides.append(slide)
    return slides


# ── PPTX 生成 ──

def _create_pptx_from_images(image_paths, output_path, slide_width=13.333, slide_height=7.5):
    """从截图列表创建图片模式 PPTX。"""
    from pptx import Presentation
    from pptx.util import Inches, Emu
    
    prs = Presentation()
    prs.slide_width = Inches(slide_width)
    prs.slide_height = Inches(slide_height)
    
    blank_layout = prs.slide_layouts[6]  # Blank layout
    
    for img_path in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        # Full-bleed image
        slide.shapes.add_picture(
            img_path,
            left=0, top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
    
    prs.save(output_path)
    return len(image_paths)


def _hex_to_rgb(hex_str):
    """#RRGGBB → (R, G, B)"""
    hex_str = hex_str.lstrip('#')
    if len(hex_str) == 3:
        hex_str = ''.join(c*2 for c in hex_str)
    return tuple(int(hex_str[i:i+2], 16) for i in (0, 2, 4))

def _create_pptx_from_structured(slides_data, output_path, style_id="03", slide_width=13.333, slide_height=7.5):
    """从结构化数据创建可编辑 PPTX。"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    
    TOKENS = {
        "03": {"bg": "FAFAF5", "fg": "333333", "accent": "2171B5", "muted": "777777",
                "font_title": "Georgia", "font_body": "Georgia"},
        "04": {"bg": "FFFFFF", "fg": "1A1A1A", "accent": "FF6B00", "muted": "888888",
                "font_title": "Calibri", "font_body": "Calibri"},
        "default": {"bg": "FFFFFF", "fg": "333333", "accent": "2171B5", "muted": "888888",
                     "font_title": "Calibri", "font_body": "Calibri"},
    }
    tk = TOKENS.get(style_id, TOKENS["default"])
    
    prs = Presentation()
    prs.slide_width = Inches(slide_width)
    prs.slide_height = Inches(slide_height)
    blank_layout = prs.slide_layouts[6]
    
    for slide_info in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        
        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(tk["bg"])
        
        # Add accent line (top)
        line = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(1.0), Inches(0.8), Inches(1.0), Pt(4)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor.from_string(tk["accent"])
        line.line.fill.background()
        
        # Add texts
        y_offset = Inches(1.2)
        for text_info in slide_info.get("texts", []):
            if not text_info["text"]:
                continue
            
            left = Inches(1.0)
            width = Inches(slide_width - 2.0)
            
            # Size mapping: HTML px → PPTX pt (approximate)
            pt_size = Pt(min(text_info["size"] * 0.7, 60))
            
            # Height proportional to font size
            if text_info["level"] == "title":
                height = Inches(1.5)
            elif text_info["level"] == "subtitle":
                height = Inches(1.0)
            else:
                height = Inches(0.5)
            
            # Color: use extracted if available, else token
            if text_info.get("color") and text_info["color"].startswith('#'):
                try:
                    rgb = _hex_to_rgb(text_info["color"])
                    color = RGBColor(*rgb)
                except:
                    color = RGBColor.from_string(tk["fg"])
            elif text_info["level"] in ("title", "subtitle"):
                color = RGBColor.from_string(tk["fg"])
            else:
                color = RGBColor.from_string(tk["muted"])
            
            # If accent color mentioned in style
            if text_info.get("color") and tk["accent"].lower() in str(text_info["color"]).lower():
                color = RGBColor.from_string(tk["accent"])
            
            txBox = slide.shapes.add_textbox(left, y_offset, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text_info["text"]
            p.font.size = pt_size
            p.font.color.rgb = color
            p.font.bold = text_info.get("bold", text_info["level"] == "title")
            p.font.name = tk["font_title"] if text_info["level"] in ("title", "subtitle") else tk["font_body"]
            
            y_offset += height + Inches(0.08)
        
        # Add note about charts
        if slide_info.get("has_chart"):
            note_box = slide.shapes.add_textbox(
                Inches(1.0), Inches(slide_height - 1.5),
                Inches(6), Inches(0.4)
            )
            tf = note_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"[图表: {slide_info.get('chart_count', 1)} 个SVG — 保真模式可完整还原]"
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor.from_string(tk["muted"])
            p.font.italic = True
    
    prs.save(output_path)
    return len(slides_data)


# ── 主导出函数 ──

def export(html_path, output_path, mode="fidelity", style_id="03", 
           serve_dir=None, width=1920, height=1080, brand_shell=None):
    """导出 HTML 幻灯片为 PPTX。

    Args:
        html_path: HTML 文件路径
        output_path: 输出 PPTX 路径
        mode: "fidelity"（保真，默认）或 "editable"（可编辑）
        style_id: 可编辑模式使用的风格 token
        serve_dir: HTTP 服务目录（如果已有可复用）
        width/height: 幻灯片像素尺寸
        brand_shell: 品牌套壳配置 {"logo_path": "..."} 或 None

    Returns:
        生成的幻灯片数量
    """
    html_path = Path(html_path)
    output_path = Path(output_path)
    
    if mode == "fidelity":
        return _export_fidelity(html_path, output_path, serve_dir, width, height, brand_shell)
    elif mode == "editable":
        n = _export_editable(html_path, output_path, style_id)
        if brand_shell:
            _apply_brand(str(output_path), str(output_path), brand=brand_shell.get("brand", "custom"))
        return n
    else:
        raise ValueError(f"Unknown mode: {mode}")


def _export_fidelity(html_path, output_path, serve_dir=None, width=1920, height=1080, brand_shell=None):
    """保真模式：截图→图片嵌入 PPTX。
    
    当 brand_shell 指定时：
    - 保持原始 16:9 视口渲染（字体不缩小、无边缘黑线）
    - 截图后用 PIL 裁剪到内容区比例，加内边距
    - 底色设为纯白，与模板底色融合
    """
    import tempfile
    from playwright.sync_api import sync_playwright
    
    # Start HTTP server if needed
    if serve_dir is None:
        serve_dir = str(html_path.parent)
    
    sys.path.insert(0, str(Path(__file__).parent))
    from srv import serve_dir as _serve, stop as _stop
    
    base_url = _serve(serve_dir)
    url = f"{base_url}/{html_path.name}"
    
    # Use Playwright for reliable screenshots
    with sync_playwright() as p:
        browser = p.chromium.launch()
        
        if brand_shell:
            # For brand template: use 2.14:1 aspect ratio (9.2" x 4.30")
            # Height 1080px → width 2311px
            screenshot_width = 2311
            page = browser.new_page(
                viewport={'width': screenshot_width, 'height': height},
                device_scale_factor=2
            )
            # Inject CSS to make slides fill the wider viewport
            page.add_init_script("""
                document.addEventListener('DOMContentLoaded', () => {
                    const style = document.createElement('style');
                    style.textContent = `
                        #dw, .s { width: 100vw !important; }
                        .s { padding: 60px 120px !important; }
                    `;
                    document.head.appendChild(style);
                });
            """)
        else:
            # Standard 16:9
            page = browser.new_page(
                viewport={'width': width, 'height': height},
                device_scale_factor=2
            )
        
        page.goto(url, wait_until='networkidle')
        page.wait_for_timeout(1500)
        
        # Count slides
        n_slides = page.evaluate("document.querySelectorAll('.s').length")
        
        if n_slides == 0:
            browser.close()
            _stop(base_url)
            raise RuntimeError("No slides found in HTML")
        
        # Screenshot each slide
        img_paths = []
        tmp_dir = tempfile.mkdtemp(prefix="html2pptx_")
        
        for i in range(n_slides):
            # Switch slide
            page.evaluate(f'''
                document.querySelectorAll('.s').forEach(s=>s.classList.remove('on'));
                document.querySelectorAll('.s')[{i}].classList.add('on');
            ''')
            page.wait_for_timeout(600)
            
            # Take screenshot
            img_path = os.path.join(tmp_dir, f"slide_{i:03d}.png")
            page.screenshot(path=img_path)
            img_paths.append(img_path)
        
        browser.close()
    
    # Build PPTX — use standard 16:9 slide size (13.333 x 7.5 inches)
    # Screenshots are 3840x2160 (2x), embedded as full-bleed on standard slides
    n = _create_pptx_from_images(img_paths, str(output_path), 13.333, 7.5)
    
    # Cleanup
    _stop(base_url)
    for p in img_paths:
        try:
            os.remove(p)
        except:
            pass
    try:
        os.rmdir(tmp_dir)
    except:
        pass
    
    # Apply brand shell if requested
    if brand_shell:
        _apply_brand(str(output_path), str(output_path), brand=brand_shell.get("brand", "custom"),
                     html_source=str(html_path))
    
    return n


def _export_editable(html_path, output_path, style_id="03"):
    """可编辑模式：解析 HTML→原生 PPTX 元素。"""
    with open(html_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    slides_data = _parse_slides(content)
    n = _create_pptx_from_structured(slides_data, str(output_path), style_id)
    return n


def _apply_brand(input_path, output_path, brand="custom", html_source=None):
    """用品牌模板母版套壳。

    封面页：填入标题/副标题，不嵌图（保留模板封面布局）
    结尾页：保留模板默认"谢谢"，不嵌图
    内容页：截图按内容区比例生成，填满无变形
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from lxml import etree

    # 定位品牌模板（用户提供自定义模板，放在 brand-template/ 目录下）
    script_dir = Path(__file__).parent
    template_path = None
    for candidate in [
        script_dir.parent / "brand-template" / "template.pptx",
        script_dir.parent.parent / "brand-template" / "assets" / "template.pptx",
    ]:
        if candidate.exists():
            template_path = str(candidate)
            break

    if not template_path:
        _apply_brand_simple(input_path, output_path)
        return
    
    src_prs = Presentation(input_path)
    tpl_prs = Presentation(template_path)
    n_slides = len(src_prs.slides)
    
    # ── 内容区精确尺寸（从模板测量） ──
    # 内容页：顶栏0.87"高，页脚线y=5.17"，内容区 9.2"×4.30"
    CONTENT_LEFT   = Inches(0.4)
    CONTENT_TOP    = Inches(0.87)
    CONTENT_WIDTH  = Inches(9.2)
    CONTENT_HEIGHT = Inches(4.30)   # 5.17 - 0.87 = 4.30"
    # 内容区宽高比 = 9.2/4.30 = 2.14:1
    
    # 删除模板自带的示例幻灯片
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
             'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}
    pres_elem = tpl_prs.part._element
    sldIdLst = pres_elem.find('p:sldIdLst', nsmap)
    slide_rids = []
    for sldId in list(sldIdLst):
        rid = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        slide_rids.append(rid)
        sldIdLst.remove(sldId)
    for rid in slide_rids:
        try:
            tpl_prs.part.drop_rel(rid)
        except Exception:
            pass
    
    # ── 从 HTML 解析标题/副标题 ──
    titles = []
    subtitles = []
    
    parsed = False
    if html_source and os.path.exists(html_source):
        with open(html_source, 'r', encoding='utf-8') as f:
            html_content = f.read()
        sections = re.split(r'<section\s+class="s[^"]*"[^>]*>', html_content)
        for section in sections[1:n_slides+1]:
            title_text = ""
            subtitle_text = ""
            for match in re.finditer(r'font-size:\s*(\d+)px[^>]*>([^<]+)<', section):
                sz = int(match.group(1))
                txt = match.group(2).strip()
                if sz >= 40 and not title_text:
                    title_text = txt
                elif 24 <= sz < 40 and not subtitle_text:
                    subtitle_text = txt
            titles.append(title_text)
            subtitles.append(subtitle_text)
        parsed = True
    
    # Ensure we have enough titles/subtitles for all slides
    while len(titles) < n_slides:
        titles.append("")
    while len(subtitles) < n_slides:
        subtitles.append("")
    
    if not parsed:
        # Fallback: 从截图的文字区域提取（不可靠，但聊胜于无）
        for slide in src_prs.slides:
            title_text = ""
            subtitle_text = ""
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                if not title_text:
                    title_text = text
                elif not subtitle_text and text != title_text:
                    subtitle_text = text
            titles.append(title_text)
            subtitles.append(subtitle_text)
    
    # ── 收集源截图（仅内容页需要） ──
    slide_images = []
    for slide in src_prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'image'):
                img_blob = shape.image.blob
                img_ext = shape.image.content_type.split('/')[-1]
                if img_ext == 'jpeg':
                    img_ext = 'jpg'
                slide_images.append((img_blob, img_ext))
                break
        else:
            slide_images.append(None)
    
    # ── 逐页套壳 ──
    layouts = tpl_prs.slide_layouts
    import tempfile
    
    for i in range(n_slides):
        # ── 封面页：填标题，不嵌图 ──
        if i == 0:
            slide = tpl_prs.slides.add_slide(layouts[0])
            # 填标题占位符 [1]
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    ph.text = titles[i][:50] if titles[i] else "演示标题"
                    for p in ph.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.name = "华文黑体"
                            run.font.size = Pt(28)
                elif ph.placeholder_format.idx == 10:
                    ph.text = subtitles[i][:60] if subtitles[i] else ""
                    for p in ph.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.name = "华文黑体"
                            run.font.size = Pt(14)
            continue
        
        # ── 结尾页：保留模板默认，不嵌图 ──
        if i == n_slides - 1:
            tpl_prs.slides.add_slide(layouts[5])
            continue
        
        # ── 内容页：截图已匹配内容区比例（2.14:1），直接嵌入 ──
        slide = tpl_prs.slides.add_slide(layouts[2])
        
        if slide_images[i] is None:
            continue
        
        img_blob, img_ext = slide_images[i]
        
        # 直接嵌入，填满内容区
        try:
            from PIL import Image
            import io
            import tempfile
            
            img = Image.open(io.BytesIO(img_blob))
            src_w, src_h = img.size
            
            # 裁剪微小空白边（检测非背景色区域）
            import numpy as np
            img_array = np.array(img)
            bg_color = np.array([250, 250, 245])  # #FAFAF5
            diff = np.abs(img_array[:,:,:3].astype(int) - bg_color)
            non_bg_mask = np.any(diff > 15, axis=2)
            
            non_bg_rows = np.where(np.any(non_bg_mask, axis=1))[0]
            non_bg_cols = np.where(np.any(non_bg_mask, axis=0))[0]
            
            if len(non_bg_rows) > 0 and len(non_bg_cols) > 0:
                content_top = max(0, non_bg_rows[0] - 2)
                content_bottom = min(src_h, non_bg_rows[-1] + 3)
                content_left = max(0, non_bg_cols[0] - 2)
                content_right = min(src_w, non_bg_cols[-1] + 3)
                img = img.crop((content_left, content_top, content_right, content_bottom))
            
            # 保存并嵌入
            tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
            img.save(tmp, 'PNG', quality=95)
            tmp.close()
            
            slide.shapes.add_picture(tmp.name, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT)
            os.unlink(tmp.name)
        except ImportError:
            # 无 PIL：直接嵌入
            tmp = tempfile.NamedTemporaryFile(suffix=f'.{img_ext}', delete=False)
            tmp.write(img_blob)
            tmp.close()
            slide.shapes.add_picture(tmp.name, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT)
            os.unlink(tmp.name)
    
    tpl_prs.save(output_path)


def _apply_brand_simple(input_path, output_path):
    """Fallback：模板不存在时，添加简单品牌条。"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    
    prs = Presentation(input_path)
    for slide in prs.slides:
        bar = slide.shapes.add_shape(
            1, 0, Inches(0),
            prs.slide_width, Inches(0.6)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(0x00, 0x50, 0x96)
        bar.line.fill.background()
        
        # 底部品牌条
        bot = slide.shapes.add_shape(
            1, 0, prs.slide_height - Inches(0.04),
            prs.slide_width, Inches(0.04)
        )
        bot.fill.solid()
        bot.fill.fore_color.rgb = RGBColor(0x00, 0x50, 0x96)
        bot.line.fill.background()
        
        credit = slide.shapes.add_textbox(
            Inches(0.3), prs.slide_height - Inches(0.35),
            Inches(3), Inches(0.2)
        )
        p = credit.text_frame.paragraphs[0]
        p.text = "Brand"
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    
    prs.save(output_path)


# ── CLI ──

if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="Export my-design HTML to PPTX")
    parser.add_argument("html", help="Input HTML file")
    parser.add_argument("output", help="Output PPTX file")
    parser.add_argument("--mode", choices=["fidelity", "editable"], default="fidelity",
                        help="Export mode (default: fidelity)")
    parser.add_argument("--style", default="03", help="Style ID for editable mode (default: 03)")
    parser.add_argument("--brand", default=None, help="Brand shell: path to logo, or 'simple' for a plain brand bar")
    args = parser.parse_args()
    
    brand = {"logo_path": args.brand} if args.brand and args.brand != "simple" else (
        {"logo_path": None} if args.brand == "simple" else None
    )
    n = export(args.html, args.output, mode=args.mode, style_id=args.style, brand_shell=brand)
    print(f"Exported {n} slides → {args.output}")
